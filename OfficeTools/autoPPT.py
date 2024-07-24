# 导入pptx库
import pptx
# 导入os库
import os

# 创建一个ppt文件
ppt = pptx.Presentation()
# 修改幻灯片的大小
ppt.slide_width = pptx.util.Inches(16)
ppt.slide_height = pptx.util.Inches(9)

# 读取下picture文件夹的所有图片
for pic in os.listdir('picture'):
    if pic != '.DS_Store':
        # 添加一个空白的幻灯片
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        # 添加一个图片
        slide.shapes.add_picture(
            'picture/' + pic, 0, 0, pptx.util.Inches(16), pptx.util.Inches(9))

# 保存ppt文件
ppt.save('test.pptx')
print('done')
